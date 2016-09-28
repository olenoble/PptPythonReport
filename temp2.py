# -*- coding: utf-8 -*-
"""
Spyder Editor

This is a temporary script file.
"""

import pptx
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
import os
os.chdir('C:/Users/Olivier/Desktop/Olivier Docs/test')

#####
def replace_in_text_frame(x):
    for p in x.text_frame.paragraphs:
        for r in p.runs:
            t = r.text
            r.text = t.replace('Insert', 'xxx')
        
def replace_in_table(x):
    for tbr in x.table.rows:
        for tbc in tbr.cells:
            ## cells can only have text ?
            #if tbc.has_text_frame:
                replace_in_text_frame(tbc)
            #elif tbc.has_table:
            #    replace_in_table(x)
            #else:
            #    pass

def is_auto_shape_to_replace(x):
    try:
        return x.auto_shape_type == MSO_SHAPE.RECTANGLE
    except:
        return False

def replace_to_graph(x, parent_slide):
    parent_slide.add_picture('plot.png', x.left, x.top, x.width, x.height)
    x.width = 0
    x.height = 0
     
def replace_in_shape(x, parent_slide):
    if is_auto_shape_to_replace(x):
        replace_to_graph(x, parent_slide)
    
    if x.has_text_frame:
        replace_in_text_frame(x)
    elif x.has_table:
        replace_in_table(x)
    else:
        print x
        #x.top = int(x.top * 1.2)
    
def replace_in_slide(x):
    for s in x.shapes:
        replace_in_shape(s, x.shapes)
        
def replace_in_presentation(x):
    for s in x.slides:
        replace_in_slide(s)
        
        #left = top = Inches(1)


prs = pptx.Presentation('template2.pptx')
replace_in_presentation(prs) 
prs.save('test2.pptx')

